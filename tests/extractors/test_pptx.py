"""Tests for the PPTX extractor functionality."""

import io
import os
from unittest.mock import MagicMock, patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.slide import Slide, Slides

from pyvisionai.extractors.pptx import ImageTask, PptxTextImageExtractor


@pytest.fixture
def test_output_dir(tmp_path):
    """Create a temporary directory for test outputs."""
    output_dir = tmp_path / "output"
    output_dir.mkdir()
    return str(output_dir)


@pytest.fixture
def mock_presentation():
    """Create a mock presentation with test slides."""
    prs = MagicMock(spec=Presentation)
    slides = MagicMock(spec=Slides)
    prs.slides = slides

    # Create test slides
    slide1 = MagicMock(spec=Slide)
    shape1 = MagicMock()
    shape1.text = "Test text 1"
    slide1.shapes = [shape1]

    slide2 = MagicMock(spec=Slide)
    shape2 = MagicMock()
    shape2.text = ""  # Empty text
    slide2.shapes = [shape2]

    # Add image relationship to slide2
    rel = MagicMock()
    rel.reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    rel.target_part.blob = b"fake_image_data"
    slide2.part.rels = {"rId1": rel}

    slides.__iter__.return_value = [slide1, slide2]
    return prs


def test_extract_text_and_images():
    """Test extracting text and images from a PPTX file."""
    mock_prs = MagicMock(spec=Presentation)
    slides = MagicMock(spec=Slides)
    mock_prs.slides = slides

    # Create test slides with text and images
    slide1 = MagicMock(spec=Slide)
    shape1, shape2 = MagicMock(), MagicMock()
    shape1.text, shape2.text = "Text 1", "Text 2"
    slide1.shapes = [shape1, shape2]

    slide2 = MagicMock(spec=Slide)
    shape3 = MagicMock()
    shape3.text = ""
    slide2.shapes = [shape3]

    # Add image relationships
    rel1, rel2 = MagicMock(), MagicMock()
    rel1.reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    rel2.reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    rel1.target_part.blob = b"image1_data"
    rel2.target_part.blob = b"image2_data"
    slide2.part.rels = {"rId1": rel1, "rId2": rel2}

    slides.__iter__.return_value = [slide1, slide2]

    with patch(
        "pyvisionai.extractors.pptx.Presentation", return_value=mock_prs
    ):
        extractor = PptxTextImageExtractor()
        texts, images = extractor.extract_text_and_images("test.pptx")

        assert len(texts) == 2
        assert texts[0] == "Text 1\nText 2"
        assert texts[1] == ""
        assert len(images) == 2
        assert images[0] == b"image1_data"
        assert images[1] == b"image2_data"


def test_extract_text_and_images_error():
    """Test error handling in text and image extraction."""
    with patch(
        "pyvisionai.extractors.pptx.Presentation",
        side_effect=PackageNotFoundError("Invalid PPTX"),
    ):
        extractor = PptxTextImageExtractor()
        with pytest.raises(PackageNotFoundError):
            extractor.extract_text_and_images("invalid.pptx")


def test_save_image(test_output_dir):
    """Test saving an image from bytes data."""
    # Create a small test image
    img = Image.new('RGB', (100, 100), color='red')
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')

    extractor = PptxTextImageExtractor()
    img_path = extractor.save_image(
        img_bytes.getvalue(), test_output_dir, "test_image"
    )

    assert os.path.exists(img_path)
    assert img_path.endswith(".jpg")

    # Verify the saved image
    saved_img = Image.open(img_path)
    assert saved_img.mode == "RGB"
    assert saved_img.size == (100, 100)


def test_save_image_error():
    """Test error handling when saving an invalid image."""
    extractor = PptxTextImageExtractor()
    with pytest.raises(Exception) as exc_info:
        extractor.save_image(
            b"invalid_image_data", "/tmp", "test_image"
        )
    assert "cannot identify image file" in str(exc_info.value)


def test_process_image_task(test_output_dir):
    """Test processing a single image task."""
    # Create a test image
    img = Image.new('RGB', (100, 100), color='blue')
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')

    task = ImageTask(
        image_data=img_bytes.getvalue(),
        image_name="test_image",
        output_dir=test_output_dir,
        index=0,
    )

    extractor = PptxTextImageExtractor()
    with patch.object(
        extractor, 'describe_image', return_value="A blue image"
    ):
        idx, description = extractor.process_image_task(task)

        assert idx == 0
        assert description == "A blue image"
        # Verify the image was cleaned up
        assert not os.path.exists(
            os.path.join(test_output_dir, "test_image.jpg")
        )


def test_process_image_task_error():
    """Test error handling in image task processing."""
    task = ImageTask(
        image_data=b"invalid_image_data",
        image_name="test_image",
        output_dir="/tmp",
        index=0,
    )

    extractor = PptxTextImageExtractor()
    with patch.object(
        extractor,
        'save_image',
        side_effect=Exception("Could not save image"),
    ):
        idx, description = extractor.process_image_task(task)

        assert idx == 0
        assert (
            "Error: Could not process image test_image" in description
        )


def test_extract_full_process(test_output_dir):
    """Test the complete extraction process."""
    mock_prs = MagicMock(spec=Presentation)
    slides = MagicMock(spec=Slides)
    mock_prs.slides = slides

    # Create test slides with text and images
    slide1 = MagicMock(spec=Slide)
    shape1 = MagicMock()
    shape1.text = "Slide 1 Text"
    slide1.shapes = [shape1]

    slide2 = MagicMock(spec=Slide)
    shape2 = MagicMock()
    shape2.text = "Slide 2 Text"
    slide2.shapes = [shape2]

    # Add image relationship
    rel = MagicMock()
    rel.reltype = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    img = Image.new('RGB', (100, 100), color='blue')
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    rel.target_part.blob = img_bytes.getvalue()
    slide2.part.rels = {"rId1": rel}

    slides.__iter__.return_value = [slide1, slide2]

    with (
        patch(
            "pyvisionai.extractors.pptx.Presentation",
            return_value=mock_prs,
        ),
        patch.object(
            PptxTextImageExtractor,
            'describe_image',
            return_value="A test image",
        ),
    ):

        extractor = PptxTextImageExtractor()
        output_path = extractor.extract("test.pptx", test_output_dir)

        assert os.path.exists(output_path)
        with open(output_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Verify markdown content
        assert "# test" in content
        assert "## Slide 1" in content
        assert "Slide 1 Text" in content
        assert "## Slide 2" in content
        assert "Slide 2 Text" in content
        assert "[Image 1]" in content
        assert "Description: A test image" in content


def test_extract_error_handling():
    """Test error handling in the main extract method."""
    extractor = PptxTextImageExtractor()
    with pytest.raises(Exception):
        extractor.extract("nonexistent.pptx", "/tmp")
